from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Turing Machine and Universal Turing Machine"
subtitle.text = "A Deep Dive into the Foundations of Computation\nYour Name\nDate"

# Slide 2: What is a Turing Machine?
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "What is a Turing Machine?"
content.text = ("- A theoretical model of computation introduced by Alan Turing in 1936\n"
                "- Consists of Tape, Head, State Register, and Table of Instructions\n"
                "- Used to understand the limits of what can be computed")

# Slide 3: Components of a Turing Machine
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Components of a Turing Machine"
content.text = ("- **Tape:** An infinitely long strip divided into cells\n"
                "- **Head:** Reads and writes symbols on the tape, moves left or right\n"
                "- **State Register:** Holds the current state of the Turing Machine\n"
                "- **Table of Instructions:** Dictates actions based on current state and read symbol")

# Slide 4: How a Turing Machine Operates
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "How a Turing Machine Operates"
content.text = ("1. Reads the symbol at the head position\n"
                "2. Looks up the corresponding instruction\n"
                "3. Writes a new symbol on the tape\n"
                "4. Moves the head left or right\n"
                "5. Transitions to a new state\n"
                "6. Continues until a halt state is reached")

# Slide 5: What is a Universal Turing Machine?
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "What is a Universal Turing Machine?"
content.text = ("- A special type of Turing Machine that can simulate any other Turing Machine\n"
                "- Reads the description of any Turing Machine and its input\n"
                "- Performs the computation of the encoded Turing Machine")

# Slide 6: How a Universal Turing Machine Works
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "How a Universal Turing Machine Works"
content.text = ("1. **Encoding:** Description of a Turing Machine and its input are encoded onto the tape\n"
                "2. **Simulation:** UTM reads the encoded information and simulates the behavior\n"
                "3. **Execution:** Follows the transition rules of the encoded Turing Machine")

# Slide 7: Significance of the Universal Turing Machine
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Significance of the Universal Turing Machine"
content.text = ("- **Theoretical Foundation:** Demonstrates that a single machine can perform any algorithmic computation\n"
                "- **Turing Completeness:** Measure of the expressive power of a system\n"
                "- **Modern Computers:** Practical implementations of the UTM concept")

# Slide 8: Applications
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Applications"
content.text = ("- **Software Simulation:** Used to simulate other computational systems\n"
                "- **Artificial Intelligence:** Principles underpin AI development\n"
                "- **Theoretical Research:** Central topic in theoretical computer science")

# Slide 9: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = ("- Turing Machine: Simple yet powerful model of computation\n"
                "- Universal Turing Machine: Versatile machine that can simulate any other Turing Machine\n"
                "- Foundations of modern computing and ongoing research")

# Slide 10: Questions
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Questions"
content.text = ("- Invite questions from the audience\n"
                "- Encourage discussion on the implications and significance of the UTM")

# Save the presentation
ppt_path = "/mnt/data/Turing_Machine_Universal_Turing_Machine_Presentation_v2.pptx"
prs.save(ppt_path)
ppt_path
